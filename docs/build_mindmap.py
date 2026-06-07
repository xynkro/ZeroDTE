#!/usr/bin/env python3
"""Build the ZeroDTE system mindmap deck (python-pptx).

3 slides: (1) hero/northstar, (2) the mindmap (northstar + 4 panels + connectors),
(3) data-flow pipeline. Palette = the system's OWN dashboard colours so it's
content-informed, not generic.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette (from frontend/index.html CSS vars) ──────────────────────────────
BG      = RGBColor(0x0B, 0x0D, 0x12)
PANEL   = RGBColor(0x15, 0x19, 0x22)
PANEL2  = RGBColor(0x11, 0x14, 0x1B)
BORDER  = RGBColor(0x2A, 0x2F, 0x3A)
TEXT    = RGBColor(0xE6, 0xE9, 0xEF)
MUTED   = RGBColor(0x8A, 0x93, 0xA6)
GREEN   = RGBColor(0x26, 0xA6, 0x9A)
RED     = RGBColor(0xEF, 0x53, 0x50)
BLUE    = RGBColor(0x5B, 0x8D, 0xEF)
ORANGE  = RGBColor(0xF0, 0xA3, 0x3C)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
GOLD    = RGBColor(0xF4, 0xC4, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def rrect(slide, x, y, w, h, fill=PANEL, line=BORDER, lw=1.0, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def connect(slide, x1, y1, x2, y2, color=BORDER, w=1.75):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    return cn


def textbox(slide, x, y, w, h, title=None, title_color=TEXT, title_size=14,
            bullets=None, body_color=TEXT, body_size=11, anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT, title_bold=True, gap_after_title=6, line_gap=3,
            body_bullet="•  "):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    first = True
    if title is not None:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = title
        r.font.bold = title_bold; r.font.size = Pt(title_size); r.font.color.rgb = title_color
        r.font.name = "Calibri"
        p.space_after = Pt(gap_after_title)
        first = False
    for b in (bullets or []):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(line_gap); p.line_spacing = 1.0
        r = p.add_run(); r.text = (body_bullet + b) if body_bullet else b
        r.font.size = Pt(body_size); r.font.color.rgb = body_color; r.font.name = "Calibri"
    return tb


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO / NORTHSTAR
# ═══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)
# accent corner motif
bar(s1, 0, 0, 0.18, 7.5, GREEN)
textbox(s1, 0.9, 1.15, 11.5, 1.2, title="ZeroDTE", title_color=TEXT, title_size=66)
textbox(s1, 0.95, 2.45, 11.5, 0.6, title="0DTE SPX Directional Credit-Spread System",
        title_color=MUTED, title_size=22, title_bold=False)
# northstar callout
ns = rrect(s1, 0.95, 3.5, 11.4, 1.7, fill=PANEL, line=GREEN, lw=1.75, radius=0.06)
bar(s1, 0.95, 3.5, 0.10, 1.7, GREEN)
textbox(s1, 1.25, 3.65, 10.9, 1.45,
        title="★  NORTHSTAR", title_color=GREEN, title_size=15,
        bullets=["The backend is the boss. Telegram, the TradingView Pine indicator, and the PWA are "
                 "parallel projections of ONE validated source of truth — never independent opinions."],
        body_color=TEXT, body_size=17, body_bullet="", gap_after_title=8)
textbox(s1, 0.95, 5.6, 11.5, 0.5,
        title="System Mindmap  ·  Principles of Operation  ·  Architecture",
        title_color=MUTED, title_size=13, title_bold=False)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE MINDMAP
# ═══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
textbox(s2, 0, 0.28, 13.333, 0.5, title="System Mindmap", title_color=TEXT, title_size=20,
        align=PP_ALIGN.CENTER)

# center node geometry
cx, cy, cw, ch = 4.95, 2.95, 3.43, 1.6   # northstar center
# panel geometry (4 corners)
pw, ph = 4.25, 2.7
TL = (0.4, 0.95);  TR = (8.69, 0.95)
BL = (0.4, 3.85);  BR = (8.69, 3.85)

# connectors first (under boxes): center corners -> panel inner corners
connect(s2, cx, cy, TL[0] + pw, TL[1] + ph, GREEN, 2)
connect(s2, cx + cw, cy, TR[0], TR[1] + ph, ORANGE, 2)
connect(s2, cx, cy + ch, BL[0] + pw, BL[1], BLUE, 2)
connect(s2, cx + cw, cy + ch, BR[0], BR[1], PURPLE, 2)

# center NORTHSTAR node
rrect(s2, cx, cy, cw, ch, fill=PANEL, line=GREEN, lw=2.0, radius=0.12)
textbox(s2, cx + 0.18, cy + 0.16, cw - 0.36, ch - 0.32,
        title="★ NORTHSTAR", title_color=GREEN, title_size=14,
        bullets=["The backend is the boss.", "All outputs are faithful projections of one validated truth."],
        body_color=TEXT, body_size=11.5, body_bullet="", line_gap=2, anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER)

# Panel A — Intent (top-left, green)
rrect(s2, *TL, pw, ph, fill=PANEL, line=BORDER); bar(s2, TL[0], TL[1], 0.09, ph, GREEN)
textbox(s2, TL[0] + 0.22, TL[1] + 0.13, pw - 0.4, ph - 0.26,
        title="🎯  INTENT — what we set out to do", title_color=GREEN, title_size=12.5,
        bullets=["Fire ONLY the backtested edge — never noise",
                 "Execute faithfully (paper fills, honest ledger)",
                 "Inform + control from anywhere, any device",
                 "Projections never diverge from the boss"],
        body_size=11, anchor=MSO_ANCHOR.MIDDLE)

# Panel B — Principles (top-right, orange)
rrect(s2, *TR, pw, ph, fill=PANEL, line=BORDER); bar(s2, TR[0], TR[1], 0.09, ph, ORANGE)
textbox(s2, TR[0] + 0.22, TR[1] + 0.13, pw - 0.4, ph - 0.26,
        title="⚙  PRINCIPLES OF OPERATION", title_color=ORANGE, title_size=12.5,
        bullets=["Backend = single source of truth",
                 "Live signal == validated backtest engine",
                 "Faithful projection (TG · Pine · PWA mirror it)",
                 "Execution truth (TG = what actually filled)",
                 "Validated edge only (+$5.5k / 5 yrs)",
                 "Resilience without noise (graceful fallback)",
                 "GitOps control (backend sealed, GitHub = bus)",
                 "Honest accounting (real fills + auto-debrief)"],
        body_size=10, line_gap=1, body_bullet="› ", anchor=MSO_ANCHOR.MIDDLE)

# Panel C — Architecture (bottom-left, blue)
rrect(s2, *BL, pw, ph, fill=PANEL, line=BORDER); bar(s2, BL[0], BL[1], 0.09, ph, BLUE)
textbox(s2, BL[0] + 0.22, BL[1] + 0.13, pw - 0.4, ph - 0.26,
        title="🧠  THE BOSS & ITS PROJECTIONS", title_color=BLUE, title_size=12.5,
        bullets=["Persistent FastAPI engine (launchd KeepAlive)",
                 "Feed → predictor → BS pricing → execution",
                 "State in live_state.json (the one truth)",
                 "3 projections: Telegram · Pine · PWA",
                 "Publisher → GitHub snapshot every 5 min"],
        body_size=11, anchor=MSO_ANCHOR.MIDDLE)

# Panel D — Integrations (bottom-right, purple)
rrect(s2, *BR, pw, ph, fill=PANEL, line=BORDER); bar(s2, BR[0], BR[1], 0.09, ph, PURPLE)
textbox(s2, BR[0] + 0.22, BR[1] + 0.13, pw - 0.4, ph - 0.26,
        title="🔌  INTEGRATIONS", title_color=PURPLE, title_size=12.5,
        bullets=["Feeds: Alpaca → IBKR → yfinance",
                 "Execution: Alpaca paper (SPY, 1/10 scale)",
                 "Pricing: Black-Scholes · CBOE chains",
                 "Macro/news: Finnhub  (not Motley Fool)",
                 "Alerts: Telegram @Tron_shaft_bot",
                 "Mirror: TradingView Pine · UI: PWA + Pages"],
        body_size=11, anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATA FLOW
# ═══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
textbox(s3, 0, 0.28, 13.333, 0.5, title="How it's put together — the flow", title_color=TEXT,
        title_size=20, align=PP_ALIGN.CENTER)

# Top pipeline: feed -> predictor -> gate -> BS -> order -> BOSS
stages = [
    ("5-min Bar Feed", "Alpaca→IBKR→yfinance", BLUE),
    ("Predictor", "run_backtest on rolling buffer", BLUE),
    ("Confluence + VWAP", "≥3 of 4 factors · gate", ORANGE),
    ("Black-Scholes", "30Δ strikes + credit", PURPLE),
    ("Alpaca paper order", "SPY · 1/10 scale", GREEN),
]
sx, sy, sw, sh, gap = 0.45, 1.25, 2.18, 1.05, 0.30
for i, (t, sub, c) in enumerate(stages):
    x = sx + i * (sw + gap)
    rrect(s3, x, sy, sw, sh, fill=PANEL, line=BORDER); bar(s3, x, sy, 0.07, sh, c)
    textbox(s3, x + 0.16, sy + 0.12, sw - 0.28, sh - 0.24, title=t, title_color=TEXT, title_size=11.5,
            bullets=[sub], body_color=MUTED, body_size=9.5, body_bullet="", gap_after_title=3)
    if i < len(stages) - 1:
        connect(s3, x + sw, sy + sh / 2, x + sw + gap, sy + sh / 2, MUTED, 1.75)

# arrow down to BOSS
boss_y = 3.05
connect(s3, sx + 4 * (sw + gap) + sw / 2, sy + sh, sx + 4 * (sw + gap) + sw / 2, boss_y, GREEN, 2)

# BOSS node (center)
bx, bw, bh = 4.6, 4.13, 1.15
rrect(s3, bx, boss_y, bw, bh, fill=PANEL, line=GREEN, lw=2.0, radius=0.12)
bar(s3, bx, boss_y, 0.1, bh, GREEN)
textbox(s3, bx + 0.25, boss_y + 0.14, bw - 0.5, bh - 0.28,
        title="THE BOSS  ·  live_state.json", title_color=GREEN, title_size=14,
        bullets=["Single source of truth — every projection derives from here"],
        body_color=MUTED, body_size=10.5, body_bullet="", anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# side inputs into the boss
rrect(s3, 0.45, 3.18, 1.9, 0.9, fill=PANEL2, line=BORDER); bar(s3, 0.45, 3.18, 0.07, 0.9, ORANGE)
textbox(s3, 0.6, 3.28, 1.7, 0.7, title="Finnhub", title_color=ORANGE, title_size=11,
        bullets=["macro · news"], body_color=MUTED, body_size=9.5, body_bullet="", gap_after_title=2)
connect(s3, 2.35, 3.63, bx, boss_y + bh / 2, MUTED, 1.5)

rrect(s3, 11.0, 3.18, 1.9, 0.9, fill=PANEL2, line=BORDER); bar(s3, 11.0, 3.18, 0.07, 0.9, PURPLE)
textbox(s3, 11.15, 3.28, 1.7, 0.7, title="CBOE", title_color=PURPLE, title_size=11,
        bullets=["IC chains"], body_color=MUTED, body_size=9.5, body_bullet="", gap_after_title=2)
connect(s3, 11.0, 3.63, bx + bw, boss_y + bh / 2, MUTED, 1.5)

# projections fanning down from the boss
proj = [
    ("Telegram", "alerts + GitOps control", GREEN),
    ("TradingView Pine", "signal mirror", BLUE),
    ("PWA — live", "localhost / Tailscale", ORANGE),
    ("GitHub Pages PWA", "snapshot · phone (read-only)", PURPLE),
]
py, pw2, ph2 = 5.35, 2.85, 1.25
pgap = 0.32
total = len(proj) * pw2 + (len(proj) - 1) * pgap
px0 = (13.333 - total) / 2
for i, (t, sub, c) in enumerate(proj):
    x = px0 + i * (pw2 + pgap)
    connect(s3, bx + bw / 2, boss_y + bh, x + pw2 / 2, py, c, 1.75)
    rrect(s3, x, py, pw2, ph2, fill=PANEL, line=BORDER); bar(s3, x, py, 0.08, ph2, c)
    textbox(s3, x + 0.18, py + 0.16, pw2 - 0.34, ph2 - 0.3, title=t, title_color=c, title_size=12.5,
            bullets=[sub], body_color=MUTED, body_size=10, body_bullet="", gap_after_title=4)

textbox(s3, 0, 7.0, 13.333, 0.4,
        title="Control loop: phone PWA → GitHub control branch → backend poller (≈60s) → applied",
        title_color=MUTED, title_size=11, title_bold=False, align=PP_ALIGN.CENTER)

out = "/Users/xynkro/Documents/Trading/ZeroDTE/docs/ZeroDTE_Mindmap.pptx"
prs.save(out)
print("saved", out)
