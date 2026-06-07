#!/usr/bin/env python3
"""Build the ZeroDTE Quant Audit & Strategy Review deck (.pptx).

Branded dark terminal aesthetic (matches the PWA): bg #0A0C10, mint-green accent,
alert red, Helvetica Neue display + Menlo mono. Content is data-driven — see
build_deck() at the bottom; the SLIDES list is populated from the swarm audit
synthesis + the chief-quant narrative.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Brand ────────────────────────────────────────────────────────────────────
BG = RGBColor(0x0A, 0x0C, 0x10)
PANEL = RGBColor(0x12, 0x16, 0x1F)
PANEL2 = RGBColor(0x18, 0x1D, 0x28)
LINE = RGBColor(0x2A, 0x30, 0x3D)
INK = RGBColor(0xE8, 0xEB, 0xF2)
INK_DIM = RGBColor(0x9A, 0xA3, 0xB4)
INK_FAINT = RGBColor(0x5D, 0x66, 0x75)
GREEN = RGBColor(0x3D, 0xDC, 0x97)
RED = RGBColor(0xFF, 0x5D, 0x6C)
BLUE = RGBColor(0x5B, 0x9D, 0xFF)
AMBER = RGBColor(0xF5, 0xB4, 0x54)
VIOLET = RGBColor(0xB5, 0x8B, 0xFF)

DISPLAY = "Helvetica Neue"
MONO = "Menlo"
BODY = "Helvetica Neue"

SEV_COLOR = {"critical": RED, "high": RED, "med": AMBER, "medium": AMBER, "low": INK_DIM}
CAT_COLOR = {
    "strategy": GREEN, "frequency": VIOLET, "timing": BLUE, "macro": AMBER,
    "backend": RED, "statistical": BLUE, "cost": AMBER,
}

EMUW, EMUH = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_autofit(tf):
    # keep text boxes from auto-resizing
    try:
        tf.word_wrap = True
    except Exception:
        pass


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = EMUW
        self.prs.slide_height = EMUH
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self, bg=BG):
        s = self.prs.slides.add_slide(self.blank)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
        _solid(r, bg)
        self.n += 1
        return s

    def rect(self, s, x, y, w, h, fill, line=None, line_w=0.75, radius=False):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        return shp

    def text(self, s, x, y, w, h, runs, size=14, color=INK, font=BODY, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=4):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        _no_autofit(tf)
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        if isinstance(runs, str):
            runs = [(runs, {})]
        first = True
        for text, opt in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = opt.get("align", align)
            p.line_spacing = opt.get("spacing", spacing)
            p.space_after = Pt(opt.get("space_after", space_after))
            p.space_before = Pt(opt.get("space_before", 0))
            for seg in (text if isinstance(text, list) else [(text, {})]):
                st, so = seg if isinstance(seg, tuple) else (seg, {})
                r = p.add_run()
                r.text = st
                r.font.size = Pt(so.get("size", size))
                r.font.name = so.get("font", font)
                r.font.bold = so.get("bold", bold)
                r.font.color.rgb = so.get("color", color)
        return tb

    def accent(self, s, x, y, color=GREEN, h=0.34):
        self.rect(s, x, y, Inches(0.07), Inches(h), color)

    def footer(self, s, label="ZeroDTE · Quant Audit"):
        self.text(s, Inches(0.55), Inches(7.04), Inches(8), Inches(0.3),
                  label, size=9, color=INK_FAINT, font=MONO)
        self.text(s, Inches(11.6), Inches(7.04), Inches(1.2), Inches(0.3),
                  f"{self.n:02d}", size=9, color=INK_FAINT, font=MONO, align=PP_ALIGN.RIGHT)

    # ── Slide templates ──────────────────────────────────────────────────────
    def title(self, kicker, title, subtitle, role):
        s = self.slide()
        # atmospheric accent rectangles
        self.rect(s, Inches(10.4), Inches(-0.6), Inches(4), Inches(3.2), None, line=LINE)
        self.text(s, Inches(0.7), Inches(1.5), Inches(11), Inches(0.4),
                  kicker, size=13, color=GREEN, font=MONO, bold=True)
        self.text(s, Inches(0.66), Inches(2.0), Inches(12), Inches(1.8),
                  title, size=46, color=INK, font=DISPLAY, bold=True, spacing=1.02)
        self.text(s, Inches(0.7), Inches(3.95), Inches(11.5), Inches(0.7),
                  subtitle, size=18, color=INK_DIM, font=BODY)
        self.rect(s, Inches(0.7), Inches(4.95), Inches(7.2), Inches(0.012), LINE)
        self.text(s, Inches(0.7), Inches(5.15), Inches(11.8), Inches(1.4),
                  role, size=13, color=INK_FAINT, font=MONO, spacing=1.3)
        self.footer(s)
        return s

    def divider(self, kicker, title, color=GREEN):
        s = self.slide()
        self.accent(s, Inches(0.7), Inches(3.0), color, h=0.5)
        self.text(s, Inches(0.95), Inches(2.92), Inches(11), Inches(0.4),
                  kicker, size=13, color=color, font=MONO, bold=True)
        self.text(s, Inches(0.92), Inches(3.35), Inches(11.8), Inches(1.2),
                  title, size=34, color=INK, font=DISPLAY, bold=True, spacing=1.03)
        self.footer(s)
        return s

    def bullets(self, title, items, accent=GREEN, kicker=None, two_col=False):
        """items: list of (head, body) or str."""
        s = self.slide()
        self.accent(s, Inches(0.55), Inches(0.62), accent)
        if kicker:
            self.text(s, Inches(0.78), Inches(0.5), Inches(11), Inches(0.3),
                      kicker, size=10.5, color=accent, font=MONO, bold=True)
        self.text(s, Inches(0.78), Inches(0.5 + (0.32 if kicker else 0)), Inches(12), Inches(0.7),
                  title, size=26, color=INK, font=DISPLAY, bold=True)
        y0 = 1.55
        cols = 2 if two_col else 1
        per = (len(items) + cols - 1) // cols
        for ci in range(cols):
            x = Inches(0.7 + ci * 6.25)
            w = Inches(5.9 if two_col else 12)
            yy = y0
            for it in items[ci * per:(ci + 1) * per]:
                head, body = it if isinstance(it, tuple) else (it, None)
                runs = [([("›  ", {"color": accent, "bold": True, "font": MONO}),
                          (head, {"bold": True, "color": INK, "size": 14.5})], {"space_after": 2})]
                if body:
                    runs.append(([("    ", {}), (body, {"color": INK_DIM, "size": 12})],
                                 {"space_after": 9, "spacing": 1.18}))
                else:
                    runs[0] = ([("›  ", {"color": accent, "bold": True, "font": MONO}),
                                (head, {"color": INK, "size": 13.5})], {"space_after": 9, "spacing": 1.2})
                self.text(s, x, Inches(yy), w, Inches(1.2), runs)
                yy += 0.55 + (0.46 if body else 0) + (0.018 * max(0, len(head) - 60))
        self.footer(s)
        return s

    def kpis(self, title, tiles, note=None, accent=GREEN, kicker=None):
        """tiles: list of (k, v, tone)."""
        s = self.slide()
        self.accent(s, Inches(0.55), Inches(0.62), accent)
        if kicker:
            self.text(s, Inches(0.78), Inches(0.5), Inches(11), Inches(0.3),
                      kicker, size=10.5, color=accent, font=MONO, bold=True)
        self.text(s, Inches(0.78), Inches(0.5 + (0.32 if kicker else 0)), Inches(12), Inches(0.7),
                  title, size=26, color=INK, font=DISPLAY, bold=True)
        n = len(tiles)
        cols = min(4, n)
        rows = (n + cols - 1) // cols
        gap = 0.28
        tw = (12.0 - gap * (cols - 1)) / cols
        th = 1.6
        for i, (k, v, tone) in enumerate(tiles):
            r, c = divmod(i, cols)
            x = Inches(0.7 + c * (tw + gap))
            y = Inches(1.7 + r * (th + 0.3))
            self.rect(s, x, y, Inches(tw), Inches(th), PANEL2, line=LINE, radius=True)
            self.rect(s, x, y, Inches(0.06), Inches(th), tone or accent)
            self.text(s, Inches(0.7 + c * (tw + gap) + 0.22), Inches(1.7 + r * (th + 0.3) + 0.22),
                      Inches(tw - 0.4), Inches(0.4), k.upper(), size=10, color=INK_FAINT, font=MONO, bold=True)
            self.text(s, Inches(0.7 + c * (tw + gap) + 0.22), Inches(1.7 + r * (th + 0.3) + 0.62),
                      Inches(tw - 0.4), Inches(0.8), v, size=27, color=tone or INK, font=MONO, bold=True)
        if note:
            self.text(s, Inches(0.7), Inches(1.7 + rows * (th + 0.3) + 0.15), Inches(12), Inches(1.2),
                      note, size=12.5, color=INK_DIM, font=BODY, spacing=1.3)
        self.footer(s)
        return s

    def flaws(self, title, cards, accent=RED, kicker="FLAW"):
        """cards: list of dict(title, severity, detail, rec). Max 4 per slide."""
        s = self.slide()
        self.accent(s, Inches(0.55), Inches(0.62), accent)
        self.text(s, Inches(0.78), Inches(0.5), Inches(11), Inches(0.3),
                  kicker, size=10.5, color=accent, font=MONO, bold=True)
        self.text(s, Inches(0.78), Inches(0.82), Inches(12), Inches(0.7),
                  title, size=24, color=INK, font=DISPLAY, bold=True)
        y = 1.78
        h = (6.6 - y) / max(1, len(cards)) - 0.18
        for c in cards:
            sev = (c.get("severity") or "med").lower()
            col = SEV_COLOR.get(sev, AMBER)
            self.rect(s, Inches(0.7), Inches(y), Inches(11.95), Inches(h), PANEL, line=LINE, radius=True)
            self.rect(s, Inches(0.7), Inches(y), Inches(0.07), Inches(h), col)
            self.text(s, Inches(0.95), Inches(y + 0.12), Inches(10.0), Inches(0.4),
                      [([(c["title"], {"bold": True, "color": INK, "size": 14.5})], {})])
            self.text(s, Inches(11.0), Inches(y + 0.14), Inches(1.5), Inches(0.3),
                      sev.upper(), size=9.5, color=col, font=MONO, bold=True, align=PP_ALIGN.RIGHT)
            body = []
            if c.get("detail"):
                body.append(([(c["detail"], {"color": INK_DIM, "size": 11.5})], {"space_after": 3, "spacing": 1.15}))
            if c.get("rec"):
                body.append(([("FIX  ", {"color": GREEN, "font": MONO, "bold": True, "size": 10}),
                              (c["rec"], {"color": INK, "size": 11.5})], {"spacing": 1.15}))
            if body:
                self.text(s, Inches(0.95), Inches(y + 0.5), Inches(11.5), Inches(h - 0.5), body)
            y += h + 0.18
        self.footer(s)
        return s

    def two_col(self, title, left_t, left_items, right_t, right_items,
                left_c=GREEN, right_c=AMBER, kicker=None):
        s = self.slide()
        self.accent(s, Inches(0.55), Inches(0.62), left_c)
        if kicker:
            self.text(s, Inches(0.78), Inches(0.5), Inches(11), Inches(0.3),
                      kicker, size=10.5, color=left_c, font=MONO, bold=True)
        self.text(s, Inches(0.78), Inches(0.5 + (0.32 if kicker else 0)), Inches(12), Inches(0.7),
                  title, size=26, color=INK, font=DISPLAY, bold=True)
        for (tt, items, cc, x) in [(left_t, left_items, left_c, 0.7), (right_t, right_items, right_c, 6.95)]:
            self.rect(s, Inches(x), Inches(1.65), Inches(5.65), Inches(5.0), PANEL, line=LINE, radius=True)
            self.rect(s, Inches(x), Inches(1.65), Inches(5.65), Inches(0.55), PANEL2)
            self.text(s, Inches(x + 0.25), Inches(1.74), Inches(5.2), Inches(0.4),
                      tt, size=13, color=cc, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            yy = 2.42
            for it in items:
                self.text(s, Inches(x + 0.25), Inches(yy), Inches(5.15), Inches(0.8),
                          [([("•  ", {"color": cc, "bold": True}), (it, {"color": INK_DIM, "size": 11.5})],
                            {"spacing": 1.18})])
                yy += 0.34 + 0.16 * (len(it) // 48)
        self.footer(s)
        return s

    def save(self, path):
        self.prs.save(path)
        return path


# ── Content ──────────────────────────────────────────────────────────────────
def build_deck(out_path):
    d = Deck()

    # 1 — Title
    d.title(
        "QUANT AUDIT · STRATEGY REVIEW · JUN 2026",
        "ZeroDTE 0DTE Engine",
        "Current architecture, the trading plan, and the flaws that matter — with a prioritized fix plan.",
        "Prepared as a senior 0DTE index-options quant. Findings are code-verified and\nadversarially reviewed. Where the edge is unproven, this report says so plainly.",
    )

    # 2 — The verdict (TL;DR)
    d.bullets("The verdict, up front", [
        ("The engine is well-built; the EVIDENCE behind it is not.",
         "Clean architecture, real gates, honest backtest. But three things quietly invalidate the 'we're validating an edge' story."),
        ("1 · You are trading the wrong side of your own edge.",
         "The validated +$5,479 is ~96% PUT-selling. Live has been ~100% CALL-selling — the side with 6 backtested trades and a structural headwind."),
        ("2 · The 'live validation' is the model grading its own homework.",
         "Every live P&L number is a Black-Scholes mid-price simulation. The broker's real fill price is never read. You are validating the model against the model."),
        ("3 · The edge is statistically marginal and cost-fragile.",
         "5-year backtest t-stat = 1.85 (below the 1.96 bar). 84% of profit came from 2022–23. It goes negative at ~$45/spread of cost in recent years."),
        ("Bottom line", "Do not scale. Do not read the green number as proof. Fix the measurement, split the book by side, and re-validate honestly."),
    ], accent=AMBER, kicker="EXECUTIVE SUMMARY")

    # 3 — Architecture: the pipeline
    d.bullets("Current architecture — how a trade is born", [
        ("Source of truth", "Persistent FastAPI backend (:8765). Feed: Alpaca SPY IEX ×10 → IBKR → yfinance. Telegram, the Pine indicator, and the PWA are parallel projections."),
        ("Signal engine", "predictor.run_backtest() rolls a 50-bar 5-min buffer through a Pine-faithful detector. A 09:30–09:45 ET observation window classifies regime (VOLATILE if obs_range/ATR > 1.5 → that day trades nothing)."),
        ("Trigger (09:45–14:00 ET)", "A Stochastic %K-vs-%D reversal cross OUT of an 80/20 extreme. Overbought cross-down → sell-call; oversold cross-up → sell-put."),
        ("8 hard gates", "confluence ≥ 3/4 · VWAP alignment · trend filter · per-side cooldown · max 3 trades/day · max 3 concurrent · VIX bucket · macro blackout · mid-session vol re-gate · 2%/day loss limit."),
        ("Pricing & exits", "Strikes at 30Δ via Black-Scholes bisection (flat vol = realized 5-min σ × 1.20). $10 SPX wing. TP at 90% of credit. No stop-ladder. Time-stop 15:30 ET."),
        ("Execution", "SPX → SPY at 1/10 scale ($1 grid). MARKET multi-leg orders to Alpaca paper."),
    ], accent=GREEN, kicker="ARCHITECTURE")

    # 4 — Config / sizing reality
    d.kpis("The live configuration (verified in .env)", [
        ("Account", "$10,000", INK), ("Risk / trade", "4% = $400", INK),
        ("Max trades/day", "3", INK), ("No new entry after", "14:00 ET", INK),
        ("Short delta", "~30Δ", INK), ("Wing", "$10 SPX", INK),
        ("Take-profit", "90% credit", GREEN), ("P&L model", "BS-mid", AMBER),
    ], note="Sizing reality: $400 budget vs ~$650 max-loss/contract → recommend_contracts floors to exactly 1 contract every trade. "
            "So the GEX size-trim and the dollar-based concurrency controls are structurally inert — they can never fire at this account size.",
       accent=GREEN, kicker="CONFIGURATION")

    # 5 — The side-selection mechanism (corrected)
    d.two_col("How the system picks CALL vs PUT  —  and why it matters",
              "THE RULE (predictor.py:388-413)", [
                  "Sell-CALL fires on an overbought cross-down, UNLESS trend == 'up'.",
                  "Sell-PUT fires on an oversold cross-up, UNLESS trend == 'down'.",
                  "Trend = EMA10−EMA30 spread vs ±0.05%. Inside that band = 'flat' → BOTH sides fire.",
                  "In a quiet up-drift the stoch hits OVERBOUGHT far more than oversold → the trigger manufactures CALL signals.",
                  "A gentle grind reads 'flat', so the trend filter never blocks those counter-trend calls.",
              ],
              "THE CONSEQUENCE", [
                  "Validated edge = sell-PUTS-on-dips (with the market's upward drift + fear premium). Robust.",
                  "Live book = sell-CALLS-on-rallies. Works only in genuine down/chop; loses in a stealth grind-up.",
                  "Backtest: 147 puts (+$5,357) vs 6 calls (+$121). The call book is essentially un-tested.",
                  "→ Live is trading the unvalidated, structurally-weak side, in the regime that punishes it.",
              ],
              left_c=BLUE, right_c=RED, kicker="THE CORE ISSUE")

    # 6 — Validation reality
    d.kpis("What the headline number actually proves", [
        ("5yr P&L", "+$5,479", GREEN), ("Per-trade", "+$35.8", INK),
        ("Per-trade σ", "$240", AMBER), ("t-stat", "1.85", RED),
        ("95% bar", "t > 1.96", INK), ("Clean trades to confirm", "~173", AMBER),
        ("Live trades so far", "~4", RED), ("From 2022–23", "84%", AMBER),
    ], note="Even five years of backtest does NOT clear the 95% significance bar. The mean edge is small relative to the noise, "
            "and the bulk of it is a 2022–23 bear-market put-selling artifact. A handful of paper trades cannot possibly confirm this.",
       accent=AMBER, kicker="STATISTICAL REALITY")

    # 7 — Divider
    d.divider("THE FLAWS", "Where it actually costs money", RED)

    # 8 — Strategy flaws
    d.flaws("Strategy — the put / call inversion", [
        {"title": "Live trades the CALL book; the validated edge is 96% PUTS",
         "severity": "critical",
         "detail": "Sell-calls-on-rallies is counter-trend in an up-drift and has only 6 backtested trades (+$121) behind it. The +$5,479 is the put book. The trend filter's 0.05% 'flat' band lets counter-trend calls through in a stealth grind.",
         "rec": "Split ALL reporting into put-book vs call-book (done). Stage a flag to suppress / down-size the call book pending a per-side backtest. Re-frame the strategy honestly as put-driven."},
        {"title": "Strikes priced on a FLAT volatility surface — no skew",
         "severity": "high",
         "detail": "tv = realized 5-min σ × 1.20 is symmetric. Real 0DTE puts carry a rich skew premium; calls are cheaper. A flat surface over-credits calls and under-credits puts — flattering exactly the side you shouldn't be selling.",
         "rec": "Add a two-parameter skew (put-IV uplift, call-IV discount) calibrated to sampled real chains; re-run the backtest before trusting any call-side attribution."},
    ], kicker="STRATEGY")

    # 9 — Cost / validation flaws
    d.flaws("Measurement — the validation is circular", [
        {"title": "P&L is a Black-Scholes mid-price simulation grading itself",
         "severity": "critical",
         "detail": "Every live pt.pnl = estimated_credit × exit_pct − $25, with both credit and exit = bs.spread_value() (theoretical mid). A full grep confirms filled_avg_price / avg_price are NEVER read. The 'shadow validation' and the backtest are two outputs of the SAME kernel — it cannot observe slippage at all.",
         "rec": "Read Alpaca filled_avg_price per leg → store broker_realized_pnl SEPARATELY from model_pnl. Judge the 4–6-week validation on broker-realized P&L only."},
        {"title": "MARKET multi-leg orders guarantee worst-case slippage",
         "severity": "high",
         "detail": "alpaca_trader.py submits type='market' mleg orders. On 0DTE wings the bid/ask is wide; market orders cross the full spread on entry AND exit. The model assumes mid on both — so realized cost is structurally worse than the $25 the backtest charges.",
         "rec": "Move to marketable-limit at mid (1 tick give), short reprice loop, market only as a time-stop fallback. Add a liquidity gate once a live NBBO feed exists."},
    ], kicker="COST / VALIDATION")

    # 10 — Backend flaws
    d.flaws("Backend logic — the live path doesn't match the backtest", [
        {"title": "Exits evaluate on the DEVELOPING 5-minute bar → phantom wins",
         "severity": "high",
         "detail": "The feed re-dispatches the still-forming bar every 60s and _check_open_wave_trades runs on every dispatch with no bar-closed guard. honest_backtest only acts on CLOSED bars (worst-first). Live can book an intra-bar TP the backtest would never see — corrupting the small sample optimistically.",
         "rec": "Gate exit evaluation to closed bars only (act on bar N when a newer timestamp arrives). Keep intra-bar dispatch for quotes/UI only."},
        {"title": "Volatility is frozen at entry — blind to intraday IV expansion",
         "severity": "high",
         "detail": "bs_realized_std is captured once and only time-to-expiry shrinks. On an afternoon news-pop the frozen-low vol under-prices the loss side: the model says −60% while you're really near +110% and breaching. It mis-prices exactly the fat-tail days.",
         "rec": "Add a conservative vol-floor ratchet (lift tv when rolling realized vol exceeds entry; never reduce a loss estimate). Re-validate."},
        {"title": "Restart / outage holes corrupt the small-sample ledger",
         "severity": "high",
         "detail": "Signal keys are in-memory only; a mid-session restart re-dispatches the day's bars and can re-book phantom duplicate trades (burning the 3/day cap, injecting fake P&L). Prior-day 0DTE positions aren't force-settled on startup before the day-gate discards them.",
         "rec": "Persist signal keys + a hard bar-age guard so backfilled signals never open trades; force-settle open prior-date 0DTE on startup; append-only ledger (done)."},
    ], kicker="BACKEND")

    # 11 — Macro / timing flaws
    d.flaws("Macro & timing — the backtest and the live system are different animals", [
        {"title": "Backtest has ZERO event awareness; live blocks event windows",
         "severity": "high",
         "detail": "honest_backtest takes every signal; live enforces a hard macro blackout. 22% of the validated exits land on FOMC/CPI/NFP/OPEX days and were MORE profitable in-sample. So live blocks the days that disproportionately built the edge — the headline number doesn't describe how it actually trades.",
         "rec": "Re-run the backtest with a macro-calendar overlay: 'event-days-included' vs 'event-days-blocked' P&L. The single most valuable missing number."},
        {"title": "Dealer-gamma (GEX) is collected but never acted on — and is inert at 1-contract sizing",
         "severity": "high",
         "detail": "Negative GEX = vol-amplified = the short strike is likelier to breach before TP. The system stamps the regime for logging and trades full size into it. The only implemented response is a size-trim that cannot fire when every trade is 1 contract.",
         "rec": "Measure breach rate on negative vs positive GEX days in the backtest first; if adverse, use GEX as a STAND-ASIDE / IC-only gate (not a size-trim)."},
        {"title": "Event-day EXIT costs unmodeled; macro gate can fail silently",
         "severity": "med",
         "detail": "Both engines charge a flat $25 on FOMC as on a quiet Tuesday; open positions exit into 3–5× wider spreads at the modeled mid. A 3× closing haircut drops the validated total to $4,629 (t→1.56); 5× to $3,779 (t→1.28). The blackout also depends on impact-string labels that can drift (e.g. a missing PCE).",
         "rec": "Add an event-day cost multiplier to the backtest; emit a morning log of exactly which events were classified high-impact and their blackout windows."},
    ], kicker="MACRO / TIMING")

    # 12 — Cost / regime curve
    d.kpis("The edge curve — where it dies", [
        ("$25 / spread", "+$5,479", GREEN), ("$35 / spread", "+$3,949", AMBER),
        ("$45 / spread", "≈ $0 *", RED), ("$60 / spread", "+$124", RED),
        ("2022–23 share", "84%", AMBER), ("Recent-regime P&L", "~$678", RED),
        ("Live breakeven (SPY)", "~$6/spr", AMBER), ("Order type", "MARKET", RED),
    ], note="* At $45/spread the 2024, 2025 and 2026 years are individually negative. Realistic 0DTE costs are $15–40 + market-order slippage on top. "
            "The edge sits INSIDE the cost-assumption error bar — and the regime you actually trade in now contributed almost none of it.",
       accent=RED, kicker="COST SENSITIVITY")

    # 13 — Recommendations
    d.bullets("Recommendations — in priority order", [
        ("1 · Confront the put/call inversion", "Report per-side. Stop reading the +$5,479 as evidence for the call book. Stage call-suppression / down-sizing pending a per-side backtest."),
        ("2 · Wire REAL fills into the ledger before any deploy decision", "Read filled_avg_price; store broker_realized_pnl beside model_pnl; judge validation on broker-realized only. Move off market orders."),
        ("3 · Fix the developing-bar exit bug", "Gate exits to closed bars so the live path mirrors the backtest it claims to validate."),
        ("4 · Run the event-overlay backtest", "Event-included vs event-blocked P&L; tag trades by GEX regime and OPEX. Cheapest high-value number."),
        ("5 · Stress cost & demand significance", "Publish the $25/$40/$50/$61 curve in every debrief. Do not scale contracts until it survives $50+ AND ~100 clean broker-realized trades."),
        ("6 · Close the restart / expiry data holes", "Bar-age guard; force-settle prior-day 0DTE on startup; append-only ledger. One lost ITM loser distorts a 4-trade sample."),
        ("7 · Price the tails honestly", "Conservative vol-floor ratchet so intraday vol expansion lifts the loss side."),
        ("8 · Act on GEX or stop claiming it as a safety control", "Validate negative-vs-positive breach rates, then use GEX as a stand-aside gate — or drop the claim."),
    ], accent=GREEN, kicker="THE FIX PLAN", two_col=True)

    # 14 — What changed vs staged
    d.two_col("What I changed tonight  vs  what I staged for your sign-off",
              "DEPLOYED NOW  (safe · no edge change)", [
                  "Put-book vs call-book split in the debrief + EOD.",
                  "Cost-sensitivity edge curve ($25→$61) in the backtest + debrief anchors.",
                  "Append-only persisted trade ledger (never date-gated).",
                  "Broker-fill instrumentation scaffold (model_pnl vs broker_realized_pnl).",
                  "Config flags added, defaulting to CURRENT behaviour.",
              ],
              "STAGED  (flag, OFF by default — your call)", [
                  "DIRECTIONAL_SUPPRESS_CALLS — kill / down-size the call book.",
                  "EXIT_ON_CLOSED_BAR_ONLY — fix the developing-bar exits.",
                  "Skew-aware pricing; vol-floor ratchet (re-validate first).",
                  "Marketable-limit execution; event-day timing + cost model.",
                  "GEX stand-aside gate; tighter trend filter.",
              ],
              left_c=GREEN, right_c=AMBER, kicker="CHANGE CONTROL")

    # 15 — Bottom line
    d.bullets("The bottom line", [
        ("You have not yet proven anything — and that's fine, if you measure honestly.",
         "The architecture is sound. The problem is the instrument panel is wired to the model, not the market."),
        ("Three numbers to live by from here:",
         "broker-realized P&L (not model P&L) · per-side edge (puts vs calls) · the cost-curve breakeven. Everything else is noise."),
        ("Do this before risking a cent of real money:",
         "≥ ~100 clean, broker-realized, closed-bar trades · positive at $50+/spread cost · put/call mix consistent with the backtest · drawdown inside −$1,581."),
        ("The discipline that makes money here is subtraction:",
         "fewer trades, the right side, real fills, and the patience to let the sample speak. The edge — if it's real — is a put-selling edge. Trade that one."),
    ], accent=BLUE, kicker="CLOSING")

    return d.save(out_path)


if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "docs/ZeroDTE_Quant_Audit.pptx"
    print("wrote", build_deck(out))
